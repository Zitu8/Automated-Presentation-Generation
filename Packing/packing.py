# Title: Image Packing in PowerPoint
# Description: This script packs a random selection of images within a specified boundary in a PowerPoint presentation.

# Importing Required Libraries
import os
import collections
import collections.abc
import random
import math
from pptx import Presentation
from pptx.util import Cm, Pt
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Constants and Parameters
LOGO_FOLDER = 'logos'
START_X_CM = 20
START_Y_CM = 12
BOUNDARY_WIDTH_CM = 15
BOUNDARY_HEIGHT_CM = 4.5
LOGO_MAX_AREA = 3
PADDING_CM = 0.1  # 0.2 cm padding between images
STEP = 0.1

# This function draws a boundary rectangle on the slide
def draw_boundary(slide, start_x, start_y, width, height):
    """Draw a boundary rectangle on the slide."""
    left = Cm(start_x)
    top = Cm(start_y)

    boundary_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, Cm(width), Cm(height)
    )
    fill = boundary_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    fill.transparency = 1.0

    line = boundary_shape.line
    line.color.rgb = RGBColor(255, 0, 0)
    line.width = Pt(2.0)

# This function selects a random number of image files from the specified folder
def select_random_images(folder, min_images=1, max_images=76):
    """Select a random number of image files from the specified folder."""
    all_image_files = [f for f in os.listdir(folder) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    return random.sample(all_image_files, min(len(all_image_files), random.randint(min_images, max_images)))

# This function resizes the images to ensure they fit within a specified maximum area
def resize_images(image_files, max_area):
    """Resize images to ensure they fit within a specified maximum area."""
    rectangles = []
    padding = PADDING_CM * 2  # account for both sides

    for img_file in image_files:
        img_path = os.path.join(LOGO_FOLDER, img_file)
        with Image.open(img_path) as img:
            width, height = img.size
            aspect_ratio = width / height
            DPI = 96

            width_cm = (width / DPI) * 2.54 - padding
            height_cm = width_cm / aspect_ratio
            area_cm = width_cm * height_cm

            if area_cm > max_area:
                area_ratio = area_cm / max_area
                side_radio = math.sqrt(area_ratio)
                height_cm /= side_radio
                width_cm /= side_radio

            rectangles.append((img_path, int(width_cm * 100), int(height_cm * 100)))

    return rectangles

# This function packs the images within the specified boundary
def pack_images(rectangles, max_width, max_height, start_x=0, start_y=0):
    """Pack the images within the specified boundary."""
    packed = []
    x = start_x + max_width
    y = start_y
    max_row_height = 0

    padding = int(PADDING_CM * 100)

    for path, w, h in sorted(rectangles, key=lambda item: (item[1], item[2]), reverse=True):
        if x - w - padding * 2 < start_x:
            x = start_x + max_width
            y += max_row_height + padding * 2
            max_row_height = 0

        if y + h + padding * 2 > start_y + max_height:
            return None

        packed.append((path, x - w - padding, y + padding, w, h))
        x -= w + padding * 2
        max_row_height = max(max_row_height, h + padding * 2)

    # Ensure equal vertical space from top and bottom
    total_img_height = y + max_row_height - start_y
    extra_space = max_height - total_img_height
    vertical_shift = extra_space // 2

    return [(path, x, y + vertical_shift, w, h) for path, x, y, w, h in packed]

# This function generates the PowerPoint presentation with the packed images and saves it
def main(i):
    image_files = select_random_images(LOGO_FOLDER)
    current_max_area = LOGO_MAX_AREA

    while current_max_area > 0:
        rectangles = resize_images(image_files, current_max_area)
        packed = pack_images(rectangles, BOUNDARY_WIDTH_CM * 100, BOUNDARY_HEIGHT_CM * 100, START_X_CM * 100, START_Y_CM * 100)
        if packed:
            break
        current_max_area -= STEP

    if not packed:
        raise ValueError("Even after resizing, the images cannot be packed within the specified boundary!")

    prs = Presentation()
    prs.slide_width = Cm(60)
    prs.slide_height = Cm(40)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Clear slide placeholders
    for shape in slide.placeholders:
        sp = shape._element
        sp.getparent().remove(sp)

    draw_boundary(slide, START_X_CM, START_Y_CM, BOUNDARY_WIDTH_CM, BOUNDARY_HEIGHT_CM)

    for path, x, y, w, h in packed:
        try:
            # Use images as they are (no need for PNG conversion)
            slide.shapes.add_picture(path, Cm(x / 100), Cm(y / 100), Cm(w / 100), Cm(h / 100))
        except Exception as e:
            print(f"Error adding image from path {path}: {e}")

    prs.save(f'packed_presentation {i} images {len(packed)}.pptx')
    print(f'packed_presentation {i}.pptx saved successfully.')


# The program starts here
if __name__ == '__main__':
    for i in range(1, 21):
        main(i)
